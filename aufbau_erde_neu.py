"""
Aufbau der Erde – neue Präsentation
Person 1: Entstehung, Erdkruste, Ob. Mantel, Unt. Mantel
Person 2: Äuß. Kern, Inn. Kern, Vulkane/Erdbeben, Bedeutung
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = "/home/user/pulsewave-app/Aufbau_der_Erde.pptx"

# ── colours ──────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0a, 0x16, 0x28)
ACCENT    = RGBColor(0x1e, 0x40, 0xaf)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x94, 0xa3, 0xb8)
GREEN     = RGBColor(0x4a, 0xde, 0x80)
BLUE      = RGBColor(0x60, 0xa5, 0xfa)
ORANGE    = RGBColor(0xfb, 0x92, 0x3c)
RED       = RGBColor(0xf8, 0x71, 0x71)
YELLOW    = RGBColor(0xfb, 0xbf, 0x24)
TEAL      = RGBColor(0x2d, 0xd4, 0xbf)
PURPLE    = RGBColor(0xc0, 0x84, 0xfc)
PINK      = RGBColor(0xf4, 0x72, 0xb6)

P1_COLOR  = TEAL      # Person 1 accent
P2_COLOR  = PURPLE    # Person 2 accent

W = Inches(13.33)
H = Inches(7.5)


# ── helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # blank
    return prs.slides.add_slide(layout)


def rect(slide, x, y, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def textbox(slide, text, x, y, w, h, size=20, bold=False,
            color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_lines(tf, lines, size=18, color=WHITE, bold=False):
    """Add multiple bullet lines to an existing text frame."""
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color


def header_bar(slide, title, accent_color, person_tag=None):
    """Dark header bar with title and optional person tag."""
    rect(slide, 0, 0, 13.33, 1.1, ACCENT)
    rect(slide, 0, 0, 0.08, 1.1, accent_color)        # left colour stripe
    textbox(slide, title, 0.25, 0.08, 10, 0.9,
            size=32, bold=True, color=WHITE)
    if person_tag:
        textbox(slide, person_tag, 11.0, 0.25, 2.1, 0.6,
                size=14, bold=True, color=accent_color, align=PP_ALIGN.RIGHT)
    # thin separator line
    rect(slide, 0, 1.1, 13.33, 0.04, accent_color)


def info_card(slide, x, y, w, h, title, lines, title_color, line_size=17):
    """Rounded card with title + bullet lines."""
    rect(slide, x, y, w, h, RGBColor(0x0f, 0x23, 0x3d))
    rect(slide, x, y, 0.06, h, title_color)            # left stripe
    textbox(slide, title, x+0.15, y+0.12, w-0.25, 0.38,
            size=15, bold=True, color=title_color)
    tb = slide.shapes.add_textbox(
        Inches(x+0.15), Inches(y+0.52), Inches(w-0.25), Inches(h-0.62))
    tb.text_frame.word_wrap = True
    add_lines(tb.text_frame, lines, size=line_size)


def bg(slide):
    rect(slide, 0, 0, 13.33, 7.5, BG)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    s = blank_slide(prs)
    bg(s)

    # decorative layer circles (cross-section feel)
    for radius, color in [
        (3.8, RGBColor(0x1e,0x3a,0x5f)),
        (3.0, RGBColor(0x1e,0x40,0xaf)),
        (2.2, RGBColor(0xfb,0x92,0x3c)),
        (1.4, RGBColor(0xf8,0x71,0x71)),
        (0.7, RGBColor(0xfb,0xbf,0x24)),
    ]:
        cx = 10.8; cy = 3.75
        r  = radius
        s.shapes.add_shape(9,   # oval
            Inches(cx-r), Inches(cy-r), Inches(r*2), Inches(r*2)
        ).fill.solid() or None
        ov = s.shapes[-1]
        ov.fill.solid(); ov.fill.fore_color.rgb = color
        ov.line.fill.background()

    # left panel
    rect(s, 0, 0, 7.5, 7.5, RGBColor(0x0a,0x16,0x28))
    rect(s, 0, 0, 0.08, 7.5, TEAL)

    textbox(s, "🌍", 0.5, 0.6, 1.5, 1.2, size=52)
    textbox(s, "Aufbau der Erde", 0.5, 1.7, 6.8, 1.5,
            size=44, bold=True, color=WHITE)
    textbox(s, "Erdkruste  ·  Erdmantel  ·  Erdkern",
            0.5, 3.1, 6.8, 0.7, size=20, color=GRAY)

    rect(s, 0.5, 3.85, 6.0, 0.04, TEAL)

    textbox(s, "👤 Person 1  —  Folien 2–5", 0.5, 4.05, 6.0, 0.5,
            size=15, color=P1_COLOR)
    textbox(s, "👤 Person 2  —  Folien 6–9", 0.5, 4.55, 6.0, 0.5,
            size=15, color=P2_COLOR)

    textbox(s, "Geographie Referat  ·  2026",
            0.5, 6.7, 6.0, 0.5, size=13, color=GRAY)


def slide_entstehung(prs):
    """Folie 2 – Entstehung der Erde (Person 1)"""
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "🌑  Entstehung der Erde", P1_COLOR, "👤 Person 1  ·  Folie 2")

    info_card(s, 0.3, 1.3, 6.1, 2.7,
        "Urknall & Sonnensystem",
        ["• Vor ~4,6 Milliarden Jahren entstand die Erde",
         "• Aus einer rotierenden Gas- und Staubscheibe",
         "• um die junge Sonne",
         "• Kleine Partikel zogen sich durch Gravitation zusammen",
         "• → immer größerer Körper: Protoerde"],
        TEAL)

    info_card(s, 6.6, 1.3, 6.4, 2.7,
        "Differenzierung (Schichtenbildung)",
        ["• Die Erde war anfangs komplett geschmolzen",
         "• Schwere Elemente (Eisen, Nickel) sanken in die Mitte",
         "• Leichtere Gesteine stiegen nach oben",
         "• → Kern, Mantel und Kruste bildeten sich"],
        TEAL)

    info_card(s, 0.3, 4.2, 6.1, 2.8,
        "Abkühlung & erste Kruste",
        ["• Oberfläche kühlte langsam ab",
         "• Erste feste Erdkruste entstand",
         "• Wasser kondensierte → erste Ozeane",
         "• Erste Atmosphäre bildete sich"],
        TEAL)

    info_card(s, 6.6, 4.2, 6.4, 2.8,
        "Wichtige Fakten",
        ["• Alter der Erde: ~4,6 Mrd. Jahre",
         "• Erdradius: 6.371 km",
         "• 5 Schichten: Kruste, ob. Mantel,",
         "  unt. Mantel, äuß. Kern, inn. Kern",
         "• Je tiefer → heißer & dichter"],
        YELLOW)


def slide_erdkruste(prs):
    """Folie 3 – Die Erdkruste (Person 1)"""
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "🟢  Die Erdkruste", GREEN, "👤 Person 1  ·  Folie 3")

    info_card(s, 0.3, 1.3, 6.1, 2.7,
        "Kontinentale Kruste",
        ["• Dicke: 30–70 km",
         "• Gestein: Granit (leicht)",
         "• Trägt die Kontinente & Gebirge",
         "• Älter und leichter als ozeanische Kruste",
         "• Temperatur: 0°C bis ~1.000°C"],
        GREEN)

    info_card(s, 6.6, 1.3, 6.4, 2.7,
        "Ozeanische Kruste",
        ["• Dicke: 5–10 km",
         "• Gestein: Basalt (schwer, dunkel)",
         "• Liegt unter den Weltmeeren",
         "• Jünger und dichter als Kontinentalkruste",
         "• Wird an Plattengrenzen neu gebildet"],
        BLUE)

    info_card(s, 0.3, 4.2, 6.1, 2.8,
        "Tektonische Platten",
        ["• Erdkruste ist in ~15 große Platten geteilt",
         "• Platten bewegen sich: 2–10 cm/Jahr",
         "• Angetrieben durch Konvektion im Mantel",
         "• An Grenzen: Erdbeben, Vulkane, Gebirge"],
        ORANGE)

    info_card(s, 6.6, 4.2, 6.4, 2.8,
        "Wichtige Fakten",
        ["• Dünnste Schicht der Erde",
         "• Nur ~1% der Erdmasse",
         "• Hauptelemente: Sauerstoff, Silizium,",
         "  Aluminium",
         "• Dichte: 2,7–3,0 g/cm³"],
        GREEN)


def slide_ob_mantel(prs):
    """Folie 4 – Oberer Erdmantel (Person 1)"""
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "🔵  Der obere Erdmantel", BLUE, "👤 Person 1  ·  Folie 4")

    info_card(s, 0.3, 1.3, 6.1, 2.7,
        "Aufbau & Eigenschaften",
        ["• Tiefe: 70 – 660 km",
         "• Zustand: zähflüssig / plastisch",
         "• Gestein: Peridotit, Olivin",
         "• Temperatur: 1.000 – 1.600°C",
         "• Dichte: 3,3 – 3,5 g/cm³"],
        BLUE)

    info_card(s, 6.6, 1.3, 6.4, 2.7,
        "Asthenosphäre",
        ["• Oberer Teil des Mantels",
         "• Plastisch verformbar wie Kaugummi",
         "• Tektonische Platten 'schwimmen' darauf",
         "• Magma kann hier entstehen",
         "• Ermöglicht Plattenbewegung"],
        TEAL)

    info_card(s, 0.3, 4.2, 6.1, 2.8,
        "Konvektion",
        ["• Heißes Gestein steigt auf",
         "• Kühlt sich ab und sinkt wieder",
         "• Wie in einem Kochtopf",
         "• Treibt die tektonischen Platten an"],
        ORANGE)

    info_card(s, 6.6, 4.2, 6.4, 2.8,
        "Wichtige Fakten",
        ["• Macht ~10% des Erdvolumens aus",
         "• Hauptelemente: Sauerstoff,",
         "  Magnesium, Silizium, Eisen",
         "• Verbindung zur Erdkruste:",
         "  Lithosphäre = Kruste + ob. Mantel"],
        BLUE)


def slide_unt_mantel(prs):
    """Folie 5 – Unterer Erdmantel (Person 1)"""
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "🟠  Der untere Erdmantel", ORANGE, "👤 Person 1  ·  Folie 5")

    info_card(s, 0.3, 1.3, 6.1, 2.7,
        "Aufbau & Eigenschaften",
        ["• Tiefe: 660 – 2.900 km",
         "• Zustand: FEST (trotz Hitze!)",
         "• Grund: extremer Druck",
         "• Temperatur: 1.600 – 3.700°C",
         "• Dichte: 3,5 – 5,6 g/cm³"],
        ORANGE)

    info_card(s, 6.6, 1.3, 6.4, 2.7,
        "Mesosphäre",
        ["• Unterer Mantel = Mesosphäre",
         "• Dickste Schicht der Erde",
         "• Gestein fest wegen Druck",
         "• Sehr langsame Konvektionsströmungen",
         "• Material: Silikate & Oxide"],
        RED)

    info_card(s, 0.3, 4.2, 6.1, 2.8,
        "Bedeutung",
        ["• Macht ~74% des Erdvolumens aus",
         "• Gesamter Mantel: ~84% Erdvolumen",
         "• Wärmeleitung vom Kern zur Kruste",
         "• Treibt langfristig Plattenbewegung"],
        ORANGE)

    info_card(s, 6.6, 4.2, 6.4, 2.8,
        "Wichtige Fakten",
        ["• Dickste Einzelschicht",
         "• Hauptgestein: Peridotit",
         "• Übergang zum Kern bei 2.900 km",
         "• Grenze heißt: Kern-Mantel-Grenze",
         "  (Gutenberg-Diskontinuität)"],
        YELLOW)


def slide_auss_kern(prs):
    """Folie 6 – Äußerer Erdkern (Person 2)"""
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "🔴  Der äußere Erdkern", RED, "👤 Person 2  ·  Folie 6")

    info_card(s, 0.3, 1.3, 6.1, 2.7,
        "Aufbau & Eigenschaften",
        ["• Tiefe: 2.900 – 5.150 km",
         "• Zustand: FLÜSSIG",
         "• Material: Eisen & Nickel",
         "• Temperatur: 4.000 – 5.000°C",
         "• Dichte: ~10–12 g/cm³"],
        RED)

    info_card(s, 6.6, 1.3, 6.4, 2.7,
        "Das Magnetfeld",
        ["• Flüssiges Eisen rotiert mit der Erde",
         "• → erzeugt elektrische Ströme",
         "• → erzeugt das Erdmagnetfeld",
         "• Magnetfeld schützt vor Sonnenwind",
         "• Ohne Magnetfeld: kein Leben auf Erde!"],
        PINK)

    info_card(s, 0.3, 4.2, 6.1, 2.8,
        "Konvektion im Kern",
        ["• Heiße Metallschmelze steigt auf",
         "• Kühlt ab und sinkt wieder",
         "• Bewegung durch Erdrotation",
         "• → Dynamo-Effekt → Magnetfeld"],
        RED)

    info_card(s, 6.6, 4.2, 6.4, 2.8,
        "Wichtige Fakten",
        ["• Dicke: ~2.250 km",
         "• Entdeckt durch Erdbebenwellen",
         "• Schallwellen können flüssige",
         "  Schicht nicht durchdringen",
         "• ~30% der Erdmasse"],
        YELLOW)


def slide_inn_kern(prs):
    """Folie 7 – Innerer Erdkern (Person 2)"""
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "🟡  Der innere Erdkern", YELLOW, "👤 Person 2  ·  Folie 7")

    info_card(s, 0.3, 1.3, 6.1, 2.7,
        "Aufbau & Eigenschaften",
        ["• Tiefe: 5.150 – 6.371 km",
         "• Zustand: FEST",
         "• Material: Eisen & Nickel",
         "• Temperatur: ~5.500°C",
         "• Radius: ~1.220 km"],
        YELLOW)

    info_card(s, 6.6, 1.3, 6.4, 2.7,
        "Warum fest trotz 5.500°C?",
        ["• Extremer Druck: 3,6 Mio. bar",
         "• Druck verhindert dass Metall schmilzt",
         "• Ähnlich wie Wasser unter Druck",
         "• = höherer Siedepunkt",
         "• Physikalisches Phänomen: Druckfestigung"],
        ORANGE)

    info_card(s, 0.3, 4.2, 6.1, 2.8,
        "Besonderheiten",
        ["• Rotiert leicht schneller als die Erde",
         "• Etwa 0,3–0,5° pro Jahr schneller",
         "• Innere Struktur noch nicht vollständig",
         "  erforscht",
         "• Heißeste Stelle der Erde"],
        YELLOW)

    info_card(s, 6.6, 4.2, 6.4, 2.8,
        "Wichtige Fakten",
        ["• ~16% des Erdvolumens (gesamter Kern)",
         "• ~32% der Erdmasse",
         "• Entdeckt 1936 von Inge Lehmann",
         "• Nachweis durch Seismologie",
         "  (Erdbebenwellen-Analyse)"],
        YELLOW)


def slide_vulkane(prs):
    """Folie 8 – Vulkane & Erdbeben (Person 2)"""
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "🌋  Vulkane & Erdbeben durch Plattenbewegung",
               ORANGE, "👤 Person 2  ·  Folie 8")

    info_card(s, 0.3, 1.3, 6.1, 2.7,
        "Wie entstehen Vulkane?",
        ["• An Plattengrenzen drückt Magma hoch",
         "• Subduzierte Platten schmelzen im Mantel",
         "• Magma steigt durch Risse auf",
         "• → Vulkane an Plattengrenzen",
         "• Beispiel: Pazifischer Feuerring"],
        RED)

    info_card(s, 6.6, 1.3, 6.4, 2.7,
        "Wie entstehen Erdbeben?",
        ["• Platten bewegen sich gegeneinander",
         "• Spannung baut sich auf",
         "• Plötzliche Entladung = Erdbeben",
         "• Stärke: Richterskala",
         "• Gefährlichste Zone: Plattengrenzen"],
        ORANGE)

    info_card(s, 0.3, 4.2, 6.1, 2.8,
        "Arten von Plattengrenzen",
        ["• Divergent: Platten entfernen sich",
         "  → Mittelozeanische Rücken, Vulkane",
         "• Konvergent: Platten stoßen zusammen",
         "  → Gebirge, Erdbeben, Vulkane",
         "• Transform: Platten gleiten aneinander",
         "  → Erdbeben (z.B. San-Andreas-Graben)"],
        TEAL)

    info_card(s, 6.6, 4.2, 6.4, 2.8,
        "Bekannte Beispiele",
        ["• Vesuv (Italien) – Konvergenzgrenze",
         "• Eyjafjallajökull (Island) – Divergenz",
         "• Erdbeben Japan 2011 – Subduktion",
         "• Himalaya – Indien trifft Eurasien",
         "• Pazifischer Feuerring – ~90% Beben"],
        RED)


def slide_bedeutung(prs):
    """Folie 9 – Bedeutung des Erdaufbaus (Person 2)"""
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "🌿  Bedeutung des Erdaufbaus für das Leben",
               GREEN, "👤 Person 2  ·  Folie 9")

    info_card(s, 0.3, 1.3, 6.1, 2.7,
        "Magnetfeld schützt das Leben",
        ["• Äußerer Kern erzeugt Magnetfeld",
         "• Magnetfeld hält Sonnenwind ab",
         "• Ohne Magnetfeld: Atmosphäre wäre",
         "  weggepustet wie bei Mars",
         "• → Kein Leben ohne Erdkern!"],
        PURPLE)

    info_card(s, 6.6, 1.3, 6.4, 2.7,
        "Wärme & Energie",
        ["• Erdinneres gibt Wärme ab",
         "• Geothermie: nutzbare Energie",
         "• Vulkane: fruchtbarer Boden",
         "• Mineralien & Rohstoffe entstehen",
         "  durch geologische Prozesse"],
        ORANGE)

    info_card(s, 0.3, 4.2, 6.1, 2.8,
        "Plattentektonik & Lebensgrundlagen",
        ["• Gebirge durch Plattenbewegung",
         "  → Regen, Flüsse, Trinkwasser",
         "• Mineralreiche Böden durch Vulkane",
         "• Kohlenstoffkreislauf durch",
         "  tektonische Prozesse reguliert",
         "• → Stabiles Klima für Leben"],
        GREEN)

    info_card(s, 6.6, 4.2, 6.4, 2.8,
        "Zusammenfassung",
        ["• Kruste  → Lebensraum & Platten",
         "• Mantel  → Konvektion & Energie",
         "• Äuß. Kern → Magnetfeld",
         "• Inn. Kern → Stabilität",
         "Merksatz: tiefer = heißer,",
         "dichter, mehr Druck"],
        TEAL)


def slide_danke(prs):
    s = blank_slide(prs)
    bg(s)

    for radius, color in [
        (3.5, RGBColor(0x0f,0x23,0x3d)),
        (2.6, RGBColor(0x1e,0x40,0xaf)),
        (1.8, RGBColor(0xfb,0x92,0x3c)),
        (1.1, RGBColor(0xf8,0x71,0x71)),
        (0.5, RGBColor(0xfb,0xbf,0x24)),
    ]:
        cx = 6.66; cy = 3.4
        s.shapes.add_shape(9,
            Inches(cx-radius), Inches(cy-radius),
            Inches(radius*2), Inches(radius*2))
        ov = s.shapes[-1]
        ov.fill.solid(); ov.fill.fore_color.rgb = color
        ov.line.fill.background()

    rect(s, 0, 0, 13.33, 7.5, RGBColor(0x0a,0x16,0x28))
    # re-draw circles on top of bg rect — swap order
    # (simpler: just use translucent bg)

    textbox(s, "🌍", 5.9, 0.3, 1.5, 1.2, size=52, align=PP_ALIGN.CENTER)
    textbox(s, "Danke fürs Zuhören!", 1.5, 1.6, 10.3, 1.3,
            size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(s, "Fragen?", 1.5, 2.85, 10.3, 0.8,
            size=28, color=GRAY, align=PP_ALIGN.CENTER)

    rect(s, 2.5, 3.75, 8.33, 0.04, TEAL)

    textbox(s, "Aufbau der Erde  ·  Erdkruste  ·  Erdmantel  ·  Erdkern",
            1.5, 3.95, 10.3, 0.6, size=15, color=GRAY, align=PP_ALIGN.CENTER)

    textbox(s, "👤 Person 1  —  Entstehung · Erdkruste · Ob. Mantel · Unt. Mantel",
            1.5, 4.65, 10.3, 0.5, size=13, color=P1_COLOR, align=PP_ALIGN.CENTER)
    textbox(s, "👤 Person 2  —  Äuß. Kern · Inn. Kern · Vulkane & Erdbeben · Bedeutung",
            1.5, 5.15, 10.3, 0.5, size=13, color=P2_COLOR, align=PP_ALIGN.CENTER)


# ── build ─────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    print("Slide 1: Titel...")
    slide_title(prs)
    print("Slide 2: Entstehung der Erde...")
    slide_entstehung(prs)
    print("Slide 3: Die Erdkruste...")
    slide_erdkruste(prs)
    print("Slide 4: Oberer Erdmantel...")
    slide_ob_mantel(prs)
    print("Slide 5: Unterer Erdmantel...")
    slide_unt_mantel(prs)
    print("Slide 6: Äußerer Erdkern...")
    slide_auss_kern(prs)
    print("Slide 7: Innerer Erdkern...")
    slide_inn_kern(prs)
    print("Slide 8: Vulkane & Erdbeben...")
    slide_vulkane(prs)
    print("Slide 9: Bedeutung...")
    slide_bedeutung(prs)
    print("Slide 10: Danke...")
    slide_danke(prs)

    prs.save(OUT)
    print(f"\nGespeichert: {OUT}")
    print(f"10 Folien — Person 1: 2-5 | Person 2: 6-9")


if __name__ == "__main__":
    main()
