from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import math

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BG    = RGBColor(0x08, 0x08, 0x08)
CARD  = RGBColor(0x11, 0x11, 0x11)
CARD2 = RGBColor(0x18, 0x18, 0x18)
YELL  = RGBColor(0xFF, 0xD6, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x88, 0x88, 0x88)
DARK  = RGBColor(0x00, 0x00, 0x00)

def blank_slide():
    layout = prs.slide_layouts[6]  # blank
    slide  = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Segoe UI'
    return txBox

def add_accent_bar(slide, t):
    add_rect(slide, 0.6, t, 0.06, 0.35, YELL)

def add_badge(slide, text, l, t):
    box = add_rect(slide, l, t, len(text)*0.12+0.3, 0.32, RGBColor(0x1a,0x1a,0x00))
    add_text(slide, text, l+0.1, t+0.04, len(text)*0.12+0.2, 0.25, 11, bold=True, color=YELL)

def add_feature_card(slide, icon, title, desc, l, t):
    add_rect(slide, l, t, 3.6, 1.6, CARD2)
    add_text(slide, icon,  l+0.18, t+0.15, 0.5,  0.4, 22)
    add_text(slide, title, l+0.18, t+0.6,  3.2,  0.35, 13, bold=True, color=WHITE)
    add_text(slide, desc,  l+0.18, t+0.98, 3.2,  0.55, 10, color=MUTED)

# ─── SLIDE 1 — Title ──────────────────────────────────────────────────────────
s = blank_slide()
add_rect(s, 0, 0, 13.33, 7.5, BG)
# Glow
glow = s.shapes.add_shape(9, Inches(7), Inches(-1), Inches(6), Inches(6))
glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x1a,0x14,0x00)
glow.line.fill.background()

add_text(s, "🎵", 1, 1.2, 1, 0.8, 48)
add_text(s, "PULSEWAVE", 1, 2.1, 8, 1.1, 72, bold=True, color=YELL)
add_text(s, "Your Music. Your Way.", 1, 3.3, 8, 0.6, 24, color=MUTED, italic=True)
add_rect(s, 1, 4.1, 1.5, 0.05, YELL)
add_text(s, "Die Musik-App der nächsten Generation", 1, 4.3, 9, 0.5, 15, color=MUTED)
add_text(s, "v1.5  ·  Windows  ·  Web  ·  iOS (bald)", 1, 6.7, 8, 0.4, 11, color=RGBColor(0x44,0x44,0x44))

# ─── SLIDE 2 — Was ist Pulsewave? ─────────────────────────────────────────────
s = blank_slide()
add_accent_bar(s, 0.7)
add_text(s, "Was ist Pulsewave?", 0.8, 0.6, 10, 0.7, 32, bold=True)
add_rect(s, 0.6, 1.5, 12.1, 0.04, RGBColor(0x22,0x22,0x22))

add_text(s, "Pulsewave ist eine kostenlose Desktop-Musik-App für Windows & Linux,\ndie Zugriff auf praktisch jeden Song der Welt bietet —\nohne Abo, ohne Registrierung, ohne Kompromisse.", 0.7, 1.7, 10, 1.4, 16, color=MUTED)

for i, (icon, label) in enumerate([
    ("🎵", "Millionen Songs"),
    ("🎛️", "7-Band Equalizer"),
    ("📱", "iOS & Web"),
    ("⭐", "Premium für €2/Mo"),
]):
    col = 0.6 + i * 3.1
    add_rect(s, col, 3.3, 2.8, 1.5, CARD2)
    add_text(s, icon,  col+0.2, 3.45, 0.6, 0.5, 26)
    add_text(s, label, col+0.2, 3.95, 2.4, 0.5, 12, bold=True, color=WHITE)

add_text(s, "Kostenlos · Open Source · Keine Werbung (mit Premium)", 0.7, 5.5, 12, 0.4, 12, color=MUTED, align=PP_ALIGN.CENTER)

# ─── SLIDE 3 — Features ───────────────────────────────────────────────────────
s = blank_slide()
add_accent_bar(s, 0.7)
add_text(s, "Features", 0.8, 0.6, 10, 0.7, 32, bold=True)
add_rect(s, 0.6, 1.5, 12.1, 0.04, RGBColor(0x22,0x22,0x22))

features = [
    ("🎵", "Riesige Bibliothek",    "Jeder Song, jeder Künstler — powered by dem weltgrößten Musikindex."),
    ("🎛️", "7-Band Equalizer",     "Sub, Bass, Mid, Treble — 6 Presets oder eigene Einstellungen."),
    ("📝", "Playlists",             "Unbegrenzt erstellen, Songs hinzufügen, umbenennen, löschen."),
    ("♥",  "Liked Songs",           "Songs mit einem Klick liken und immer griffbereit haben."),
    ("🔀", "Shuffle & Repeat",      "Volle Queue-Kontrolle mit Shuffle, Repeat All und Repeat One."),
    ("⏰", "Sleep Timer (Premium)", "Musik stoppt automatisch nach gewählter Zeit."),
]
for i, (icon, title, desc) in enumerate(features):
    col = 0.6 + (i % 3) * 4.05
    row = 1.7 + (i // 3) * 1.85
    add_feature_card(s, icon, title, desc, col, row)

# ─── SLIDE 4 — Premium ────────────────────────────────────────────────────────
s = blank_slide()
add_accent_bar(s, 0.7)
add_text(s, "Premium — Nur €2/Monat", 0.8, 0.6, 11, 0.7, 32, bold=True)
add_rect(s, 0.6, 1.5, 12.1, 0.04, RGBColor(0x22,0x22,0x22))

# Free card
add_rect(s, 0.6, 1.7, 5.5, 5.3, CARD2)
add_text(s, "Free", 0.85, 1.9, 4, 0.5, 20, bold=True, color=MUTED)
add_text(s, "€0", 0.85, 2.4, 4, 0.8, 44, bold=True, color=WHITE)
add_text(s, "/Monat", 2.0, 2.65, 3, 0.4, 13, color=MUTED)
for item in ["✓  Unbegrenztes Streaming", "✓  Equalizer", "✓  Playlists & Liked Songs", "✓  Suchfunktion", "✗  Werbung vorhanden", "✗  Kein Sleep Timer"]:
    col = RGBColor(0x4c,0xaf,0x50) if item.startswith("✓") else RGBColor(0x55,0x55,0x55)
    add_text(s, item, 0.85, 3.4 + ["✓  Unbegrenztes Streaming","✓  Equalizer","✓  Playlists & Liked Songs","✓  Suchfunktion","✗  Werbung vorhanden","✗  Kein Sleep Timer"].index(item)*0.42, 5, 0.38, 12, color=col)

# Premium card
add_rect(s, 6.8, 1.5, 0.04, 5.5, YELL)
add_rect(s, 7.1, 1.7, 5.6, 5.3, RGBColor(0x14,0x12,0x00))
add_text(s, "⭐ BELIEBT", 7.3, 1.75, 4, 0.35, 10, bold=True, color=YELL)
add_text(s, "Premium", 7.3, 2.1, 4, 0.5, 20, bold=True, color=YELL)
add_text(s, "€2", 7.3, 2.55, 4, 0.8, 44, bold=True, color=YELL)
add_text(s, "/Monat", 8.45, 2.8, 3, 0.4, 13, color=MUTED)
for item in ["✓  Alles aus Free", "✓  Keine Werbung", "✓  Sleep Timer", "✓  Crossfade", "✓  Gold Theme", "✓  Priority Support"]:
    add_text(s, item, 7.3, 3.4 + ["✓  Alles aus Free","✓  Keine Werbung","✓  Sleep Timer","✓  Crossfade","✓  Gold Theme","✓  Priority Support"].index(item)*0.42, 5, 0.38, 12, color=YELL)

# ─── SLIDE 5 — Plattformen / Download ─────────────────────────────────────────
s = blank_slide()
add_accent_bar(s, 0.7)
add_text(s, "Verfügbar auf", 0.8, 0.6, 10, 0.7, 32, bold=True)
add_rect(s, 0.6, 1.5, 12.1, 0.04, RGBColor(0x22,0x22,0x22))

platforms = [
    ("🪟", "Windows",   "Windows 10 / 11\nx64 · Auto-Updates\nv1.5 verfügbar"),
    ("🐧", "Linux",     "Ubuntu 20.04+\nAppImage · x64\nDemnächst"),
    ("🌐", "Web / iOS", "Läuft im Browser\nPWA installierbar\nSafari-Support"),
    ("🍎", "App Store", "iOS & iPadOS\nNative App\nBald verfügbar"),
]
for i, (icon, name, desc) in enumerate(platforms):
    col = 0.6 + i * 3.1
    add_rect(s, col, 1.8, 2.85, 3.8, CARD2)
    add_text(s, icon, col+0.2, 2.0, 0.8, 0.7, 36)
    add_text(s, name, col+0.2, 2.75, 2.4, 0.45, 16, bold=True)
    add_text(s, desc, col+0.2, 3.25, 2.4, 1.4, 11, color=MUTED)
    if name in ("Web / iOS", "App Store"):
        add_badge(s, "COMING SOON", col+0.2, 4.9)
    else:
        add_badge(s, "VERFÜGBAR", col+0.2, 4.9)

add_text(s, "welias123.github.io/pulsewave-app", 0.6, 6.7, 12.1, 0.4, 12, color=MUTED, align=PP_ALIGN.CENTER)

# ─── SLIDE 6 — Roadmap ────────────────────────────────────────────────────────
s = blank_slide()
add_accent_bar(s, 0.7)
add_text(s, "Roadmap", 0.8, 0.6, 10, 0.7, 32, bold=True)
add_rect(s, 0.6, 1.5, 12.1, 0.04, RGBColor(0x22,0x22,0x22))

steps = [
    ("✅", "v1.3", "Lazy Loading, 35 Genres, 700 Songs"),
    ("✅", "v1.5", "Discord RPC, Mini Player, Passwort-Hashing"),
    ("🔄", "v1.6", "Linux AppImage Release"),
    ("🔜", "v2.0", "iOS Web-App & PWA"),
    ("🍎", "v3.0", "Offizieller App Store Release"),
]
for i, (icon, ver, desc) in enumerate(steps):
    t = 1.8 + i * 0.98
    done = icon == "✅"
    add_rect(s, 0.6, t, 12.0, 0.82, CARD2 if not done else RGBColor(0x0d,0x12,0x00))
    add_text(s, icon, 0.8, t+0.18, 0.5, 0.5, 18)
    add_text(s, ver,  1.4, t+0.18, 1.0, 0.45, 14, bold=True, color=YELL)
    add_text(s, desc, 2.5, t+0.18, 9.5, 0.45, 13, color=WHITE if done else MUTED)
    if i < 4:
        add_rect(s, 0.9, t+0.82, 0.04, 0.16, RGBColor(0x33,0x33,0x33))

# ─── SLIDE 7 — Closing ────────────────────────────────────────────────────────
s = blank_slide()
add_rect(s, 0, 0, 13.33, 7.5, BG)
glow2 = s.shapes.add_shape(9, Inches(-1), Inches(2), Inches(6), Inches(6))
glow2.fill.solid(); glow2.fill.fore_color.rgb = RGBColor(0x14,0x11,0x00)
glow2.line.fill.background()

add_text(s, "🎵", 5.9, 0.9, 1.5, 1.0, 48, align=PP_ALIGN.CENTER)
add_text(s, "Pulsewave", 2, 1.9, 9.33, 1.2, 56, bold=True, color=YELL, align=PP_ALIGN.CENTER)
add_text(s, "Your Music. Your Way.", 2, 3.2, 9.33, 0.6, 20, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
add_rect(s, 4.5, 4.1, 4.33, 0.04, YELL)
add_text(s, "welias123.github.io/pulsewave-app", 2, 4.3, 9.33, 0.5, 13, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, "Kostenlos downloaden · Premium für €2/Monat · Bald im App Store", 1.5, 5.1, 10.33, 0.5, 12, color=RGBColor(0x44,0x44,0x44), align=PP_ALIGN.CENTER)

prs.save('/home/user/pulsewave-app/Pulsewave_Präsentation.pptx')
print("Fertig!")
