"""
Generates Earth-layer diagrams and adds them to Aufbau_der_Erde.pptx
"""

import io
import math
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyArrowPatch, Arc
import matplotlib.patheffects as pe
from pptx import Presentation
from pptx.util import Inches

PPTX_PATH = "/home/user/pulsewave-app/Aufbau_der_Erde.pptx"

BG = "#0a1628"
COLORS = {
    "kruste":  "#4ade80",
    "omantel": "#60a5fa",
    "umantel": "#fb923c",
    "akern":   "#f87171",
    "ikern":   "#fbbf24",
}


# ── helpers ──────────────────────────────────────────────────────────────────

def fig_to_buf(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def add_image(slide, buf, left_in, top_in, width_in):
    pic = slide.shapes.add_picture(buf, Inches(left_in), Inches(top_in),
                                   width=Inches(width_in))
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


# ── diagram functions ─────────────────────────────────────────────────────────

def make_earth_sphere():
    """Slide 1 & 7 – glowing blue Earth ball."""
    fig, ax = plt.subplots(figsize=(4, 4), facecolor=BG)
    ax.set_facecolor(BG)
    ax.set_xlim(-1.3, 1.3); ax.set_ylim(-1.3, 1.3)
    ax.set_aspect("equal"); ax.axis("off")

    # glow rings
    for r, alpha in [(1.25, 0.04), (1.18, 0.07), (1.10, 0.11)]:
        ax.add_patch(mpatches.Circle((0, 0), r, color="#38bdf8", alpha=alpha, zorder=1))

    # gradient sphere via stacked transparent circles
    N = 80
    for i in range(N, 0, -1):
        r = i / N
        c = plt.cm.cool(1 - r * 0.6)
        ax.add_patch(mpatches.Circle((0, 0), r, color=c, alpha=0.04, zorder=2))

    # main circle
    ax.add_patch(mpatches.Circle((0, 0), 1.0, color="#1e40af", zorder=2))

    # continent blobs (rough shapes)
    blobs = [
        ([-0.35, -0.25, -0.1,  0.05,  0.0, -0.2, -0.35],
         [ 0.30,  0.55,  0.60, 0.45,  0.20, 0.15,  0.30]),
        ([ 0.10,  0.35,  0.45,  0.25,  0.10],
         [ 0.25,  0.20,  0.50,  0.60,  0.25]),
        ([-0.55, -0.30, -0.30, -0.55, -0.55],
         [-0.20, -0.15, -0.60, -0.55, -0.20]),
        ([ 0.10,  0.45,  0.40,  0.10,  0.10],
         [-0.30, -0.25, -0.70, -0.65, -0.30]),
    ]
    for xs, ys in blobs:
        ax.fill(xs, ys, color="#166534", alpha=0.75, zorder=3)

    # clip all drawing to the circle
    circle_clip = mpatches.Circle((0, 0), 1.0, transform=ax.transData)
    for child in ax.get_children():
        try:
            child.set_clip_path(circle_clip)
        except Exception:
            pass

    ax.set_title("🌍 Erde", color="white", fontsize=13, pad=4)
    return fig_to_buf(fig)


def make_cross_section():
    """Slide 2 – Earth cross-section with labeled layers."""
    fig, ax = plt.subplots(figsize=(4.5, 4.5), facecolor=BG)
    ax.set_facecolor(BG)
    ax.set_xlim(-1.35, 1.35); ax.set_ylim(-1.35, 1.35)
    ax.set_aspect("equal"); ax.axis("off")

    layers = [
        (1.00, COLORS["ikern"],   "Inn. Kern\n5 150–6 371 km"),
        (0.80, COLORS["akern"],   "Äuß. Kern\n2 900–5 150 km"),
        (0.58, COLORS["umantel"], "Unt. Mantel\n660–2 900 km"),
        (0.40, COLORS["omantel"], "Ob. Mantel\n70–660 km"),
        (0.26, COLORS["kruste"],  "Erdkruste\n0–70 km"),
    ]
    for radius, color, label in layers:
        ax.add_patch(mpatches.Circle((0, 0), radius, color=color, zorder=2))

    # cut-away: white half to show layers
    theta = np.linspace(0, np.pi, 300)
    xs = np.concatenate([[0], np.cos(theta) * 1.05, [0]])
    ys = np.concatenate([[0], np.sin(theta) * 1.05, [0]])
    ax.fill(xs, ys, color=BG, zorder=3)

    # dividing line
    ax.plot([0, 0], [-1.05, 1.05], color="white", lw=0.5, alpha=0.4, zorder=4)

    # labels with leader lines (left side)
    label_data = [
        (0.00,  0.13, "Kruste",    COLORS["kruste"]),
        (-0.33, 0.10, "Ob. Mantel",COLORS["omantel"]),
        (-0.49, 0.0,  "Unt. Mantel",COLORS["umantel"]),
        (-0.69, 0.0,  "Äuß. Kern", COLORS["akern"]),
        (-0.88, 0.0,  "Inn. Kern", COLORS["ikern"]),
    ]
    for r_mid, angle_off, lbl, col in label_data:
        ax.annotate(lbl, xy=(r_mid * math.cos(math.pi * 0.75 + angle_off),
                              r_mid * math.sin(math.pi * 0.75 + angle_off)),
                    xytext=(-1.25, r_mid * 0.45 + angle_off),
                    color=col, fontsize=7, fontweight="bold",
                    arrowprops=dict(arrowstyle="-", color=col, lw=0.8),
                    zorder=5)

    ax.set_title("Schichten der Erde", color="white", fontsize=11, pad=4)
    return fig_to_buf(fig)


def make_plates():
    """Slide 3 – simplified tectonic plates map."""
    fig, ax = plt.subplots(figsize=(5, 3), facecolor=BG)
    ax.set_facecolor("#0d2137")
    ax.set_xlim(0, 10); ax.set_ylim(0, 6)
    ax.axis("off")

    # ocean background
    ax.add_patch(mpatches.FancyBboxPatch((0, 0), 10, 6,
                 boxstyle="round,pad=0.1", color="#1e3a5f", zorder=0))

    # rough plate outlines (simplified polygons)
    plates = [
        # (polygon_points, label, color)
        ([(0.5,3),(2,5.5),(4,5.5),(5,3.5),(3,2),(1,2),(0.5,3)],
         "Nordamerik.", "#166534"),
        ([(5,3.5),(7,5.5),(9,5.5),(9.5,3),(7.5,2.5),(5,3.5)],
         "Eurasisch", "#1e40af"),
        ([(2,2),(3,2),(5,1),(5,-0.5),(1,-0.5),(2,2)],
         "Südamerik.", "#7c3aed"),
        ([(5,-0.5),(5,1),(7.5,2.5),(9.5,0.5),(9.5,-0.5),(5,-0.5)],
         "Afrikan.", "#b45309"),
        ([(7.5,2.5),(9.5,3),(9.5,0.5),(7.5,2.5)],
         "Indisch-Austr.", "#0f766e"),
    ]
    for poly, lbl, col in plates:
        xs = [p[0] for p in poly]
        ys = [p[1] for p in poly]
        ax.fill(xs, ys, color=col, alpha=0.55, zorder=1)
        ax.plot(xs + [xs[0]], ys + [ys[0]], color="white", lw=0.8, alpha=0.6, zorder=2)
        cx = sum(xs) / len(xs)
        cy = sum(ys) / len(ys)
        ax.text(cx, cy, lbl, color="white", fontsize=6, ha="center", va="center",
                fontweight="bold", zorder=3)

    # plate movement arrows
    arrows = [(4.8, 3.4, -0.3, 0.2), (5.2, 3.4, 0.3, 0.2),
              (5.0, 1.0, 0.0, -0.4)]
    for x, y, dx, dy in arrows:
        ax.annotate("", xy=(x+dx, y+dy), xytext=(x, y),
                    arrowprops=dict(arrowstyle="->", color="#fbbf24",
                                   lw=1.2), zorder=4)

    ax.set_title("Tektonische Platten", color="white", fontsize=11, pad=4)
    return fig_to_buf(fig)


def make_convection():
    """Slide 4 – mantle convection diagram."""
    fig, ax = plt.subplots(figsize=(4.5, 3.5), facecolor=BG)
    ax.set_facecolor(BG)
    ax.set_xlim(0, 10); ax.set_ylim(0, 7)
    ax.axis("off")

    # layers
    ax.add_patch(mpatches.FancyBboxPatch((0, 5.5), 10, 1.4,
                 boxstyle="square", color=COLORS["kruste"], alpha=0.7, zorder=1))
    ax.add_patch(mpatches.FancyBboxPatch((0, 1.5), 10, 4.0,
                 boxstyle="square", color=COLORS["umantel"], alpha=0.4, zorder=1))
    ax.add_patch(mpatches.FancyBboxPatch((0, 0.0), 10, 1.5,
                 boxstyle="square", color=COLORS["akern"], alpha=0.6, zorder=1))

    ax.text(0.3, 6.1, "Erdkruste", color=BG, fontsize=8, fontweight="bold", zorder=2)
    ax.text(0.3, 0.5, "Äuß. Kern (flüssig)", color=BG, fontsize=7, fontweight="bold", zorder=2)
    ax.text(0.3, 3.4, "Erdmantel", color="white", fontsize=9, fontweight="bold",
            alpha=0.9, zorder=2)

    # convection arrows (two cells)
    def conv_cell(cx, direction=1):
        # rising hot
        ax.annotate("", xy=(cx, 5.3), xytext=(cx, 1.7),
                    arrowprops=dict(arrowstyle="->", color="#ef4444",
                                   lw=2, connectionstyle="arc3,rad=0"))
        # top flow
        ax.annotate("", xy=(cx + direction*2.2, 5.1), xytext=(cx + 0.1*direction, 5.1),
                    arrowprops=dict(arrowstyle="->", color="#fb923c", lw=1.5))
        # sinking cool
        ax.annotate("", xy=(cx + direction*2.2, 1.7), xytext=(cx + direction*2.2, 5.1),
                    arrowprops=dict(arrowstyle="->", color="#60a5fa",
                                   lw=2, connectionstyle="arc3,rad=0"))
        # bottom flow
        ax.annotate("", xy=(cx + 0.1*direction, 1.9), xytext=(cx + direction*2.1, 1.9),
                    arrowprops=dict(arrowstyle="->", color="#93c5fd", lw=1.5))

    conv_cell(2.5, direction=1)
    conv_cell(7.5, direction=-1)

    ax.text(3.5, 3.5, "🔴 heiß\nsteigt auf", color="#fca5a5", fontsize=7,
            ha="center", va="center", zorder=3)
    ax.text(6.5, 3.5, "🔵 kalt\nsinkt", color="#93c5fd", fontsize=7,
            ha="center", va="center", zorder=3)

    ax.set_title("Konvektion im Erdmantel", color="white", fontsize=11, pad=4)
    return fig_to_buf(fig)


def make_core_magnetic():
    """Slide 5 – Earth core + magnetic field lines."""
    fig, ax = plt.subplots(figsize=(4, 4.5), facecolor=BG)
    ax.set_facecolor(BG)
    ax.set_xlim(-2.5, 2.5); ax.set_ylim(-3.0, 3.0)
    ax.set_aspect("equal"); ax.axis("off")

    # outer core
    ax.add_patch(mpatches.Circle((0, 0), 1.5, color=COLORS["akern"],
                                  alpha=0.55, zorder=2))
    # inner core
    ax.add_patch(mpatches.Circle((0, 0), 0.8, color=COLORS["ikern"],
                                  alpha=0.85, zorder=3))

    ax.text(0, 0, "Inn.\nKern", color=BG, fontsize=8, ha="center",
            va="center", fontweight="bold", zorder=4)
    ax.text(1.1, 0.35, "Äuß.\nKern", color="white", fontsize=7,
            ha="center", va="center", zorder=4)

    # magnetic field lines (curved arcs)
    theta = np.linspace(0, np.pi, 120)
    for scale, alpha in [(2.0, 0.55), (2.5, 0.35), (1.6, 0.4)]:
        xs_r =  scale * np.sin(theta)
        ys_r =  scale * np.cos(theta) - 0.0
        ax.plot(xs_r, ys_r, color="#38bdf8", lw=1.2, alpha=alpha, zorder=1)
        ax.plot(-xs_r, ys_r, color="#38bdf8", lw=1.2, alpha=alpha, zorder=1)

    # poles
    ax.annotate("N", xy=(0, 2.55), color="#7dd3fc", fontsize=12,
                fontweight="bold", ha="center")
    ax.annotate("S", xy=(0, -2.65), color="#7dd3fc", fontsize=12,
                fontweight="bold", ha="center")

    ax.set_title("Erdkern + Magnetfeld", color="white", fontsize=11, pad=4)
    return fig_to_buf(fig)


# ── main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation(PPTX_PATH)

    generators = [
        (0, make_earth_sphere,    6.0, 0.9, 3.8),   # Titel
        (1, make_cross_section,   6.0, 0.9, 3.8),   # Überblick
        (2, make_plates,          5.8, 1.5, 4.0),   # Erdkruste
        (3, make_convection,      5.7, 1.5, 4.0),   # Erdmantel
        (4, make_core_magnetic,   6.1, 0.8, 3.6),   # Erdkern
        (6, make_earth_sphere,    3.5, 1.5, 3.0),   # Danke (centered)
    ]

    for slide_idx, gen_fn, left, top, width in generators:
        slide = prs.slides[slide_idx]
        print(f"Slide {slide_idx+1}: generating {gen_fn.__name__}...")
        buf = gen_fn()
        add_image(slide, buf, left, top, width)
        print(f"  Done.")

    prs.save(PPTX_PATH)
    print(f"\nSaved: {PPTX_PATH}")


if __name__ == "__main__":
    main()
